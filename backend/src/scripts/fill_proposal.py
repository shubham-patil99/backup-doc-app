#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
fill_proposal.py  --  HPE Proposal generator
=============================================

Template: hpe-proposal-template.pptx  (3 slides)
  Slide 0  "Cover 1"           -- {{docTitle}} {{subtitle}} {{date}} {{opeId}}
  Slide 1  "Title and Content" -- {{breadcrumb}} / {{moduleNum}} {{moduleName}} / {{moduleBody}}
  Slide 2  "Thank You"         -- closing slide

Typography spec (industry-standard):
  {{breadcrumb}}  → 30pt, Bold, HPE Green
  {{moduleName}}  → 18pt, Bold, Dark
  {{moduleBody}}  → 18pt, Regular, Dark
  Body content    →  18pt, Regular, Dark  (textboxes)
  Table header    →  18pt, Bold, White on HPE Green
  Table data      →  18pt, Regular, Dark

Key architecture decision:
  Body content is NEVER written into the placeholder (idx=1).
  The placeholder is always cleared and removed after header tokens are filled.
  ALL body text and tables are rendered as explicitly positioned textboxes/tables
  so that spacing, indent, and font are 100% consistent across every slide.
"""

import sys, json, copy, html as html_mod, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

try:
    from bs4 import BeautifulSoup, NavigableString
    HAS_BS4 = True
except ImportError:
    HAS_BS4 = False

PML = "http://schemas.openxmlformats.org/presentationml/2006/main"
REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# ── Design constants ──────────────────────────────────────────────────────────
HPE_GREEN  = RGBColor(0x00, 0x60, 0x35)
HPE_GREEN_LIGHT = RGBColor(0xE8, 0xF5, 0xED)   # subtle section tint (unused for now)
STRIPE_ODD = RGBColor(0xF2, 0xF2, 0xF2)
STRIPE_EVN = RGBColor(0xFF, 0xFF, 0xFF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK  = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_MID   = RGBColor(0x44, 0x44, 0x44)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)

# ── Header typography constants ───────────────────────────────────────────────
# These apply to the template placeholder text (filled via fill_shape / replace_in_paragraph)
FONT_BREADCRUMB = Pt(30)   # {{breadcrumb}} — 30pt Bold
FONT_MODULE_NAME = Pt(18)  # {{moduleName}} — 18pt Bold
FONT_MODULE_BODY = Pt(18)  # {{moduleBody}} — 18pt Regular

# ── Layout geometry ───────────────────────────────────────────────────────────
# Header zone: breadcrumb (30pt) + module name (18pt) ≈ 1.30" total
# Breadcrumb at ~0.15", module name at ~0.65", body content starts at ~1.30"
CONTENT_LEFT   = Inches(0.42)
CONTENT_TOP    = Inches(1.30)   # body textboxes start here (pushed down for larger header)
CONTENT_WIDTH  = Inches(12.49)
CONTENT_BOTTOM = Inches(6.90)   # safe bottom margin above footer

# Table geometry — row heights sized for 18pt font + cell padding
TABLE_LEFT  = Inches(0.42)
TABLE_W     = Inches(12.49)
TABLE_ROW_H = Emu(int(Inches(0.45)))   # 18pt + ~12pt padding = 0.417\" → 0.45\" safe
TABLE_HDR_H = Emu(int(Inches(0.50)))   # header row slightly taller for weight

# Font sizes — body content textboxes (18pt matches moduleName/moduleBody)
FONT_BODY  = Pt(18)
FONT_TABLE = Pt(18)
FONT_HDR   = Pt(18)

# ── Pagination constants ──────────────────────────────────────────────────────
# Derived from 18pt / 115% line spacing:
#   18pt / 72 = 0.250\"  ×  1.15 = 0.2875\"  +5% render buffer → 0.302\"/line
# Available body area: 6.90\" - 1.30\" = 5.60\"
# Usable after gaps/padding: 5.30\"  →  floor(5.30 / 0.302) = 17 lines max
MAX_LINES_PER_SLIDE      = 17
# Table: 10 data rows × 0.45\" + 0.50\" header = 5.00\" ≤ 5.60\" available
MAX_TABLE_ROWS_PER_SLIDE = 10

# Rendering metrics at 18pt / 115% line spacing
LINE_H   = Inches(0.302)   # one rendered line height (18pt × 1.15 × 1.05 safety)
TEXT_GAP = Inches(0.08)    # vertical gap between stacked text blocks

# Characters per line at 18pt across 12.49" content width
# avg char width at 18pt ≈ 0.125\" → 12.49 / 0.125 = ~100, use 90 (10% margin)
CHARS_PER_LINE = 90

# ── List indent geometry — scaled for 18pt readability ───────────────────────
LIST_HANG   = int(Inches(0.28))   # hanging indent (wider for 18pt bullet)
LIST_BASE_L = int(Inches(0.40))   # left margin for depth-0 list items
LIST_DEPTH  = int(Inches(0.25))   # extra indent per nesting level


# ─── Placeholder font applicator ──────────────────────────────────────────────

def _apply_run_font(run, size_pt: Pt, bold: bool = False,
                    color: RGBColor = None, italic: bool = False):
    """Apply font properties to a python-pptx Run object."""
    run.font.size  = size_pt
    run.font.bold  = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def _apply_para_font(para, size_pt: Pt, bold: bool = False,
                     color: RGBColor = None, italic: bool = False):
    """Apply font properties to every run in a paragraph."""
    for run in para.runs:
        _apply_run_font(run, size_pt, bold, color, italic)
    # Also set directly on the paragraph's default run properties for safety
    try:
        rPr = para._p.find(qn("a:r"))
        if rPr is None:
            return
    except Exception:
        pass


def _style_header_placeholder(slide, fields: dict):
    """
    After token replacement, walk each placeholder and enforce typography:
      - breadcrumb text  →  30pt Bold HPE Green
      - moduleName text  →  18pt Bold Dark
      - moduleBody text  →  18pt Regular Dark  (LAYOUT_CONTENT only)
    This works at the shape level, inspecting the replaced text content.
    """
    breadcrumb_val  = fields.get("{{breadcrumb}}", "")
    module_name_val = fields.get("{{moduleName}}", "")
    module_body_val = fields.get("{{moduleBody}}", "")

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        for para in tf.paragraphs:
            text = "".join(r.text for r in para.runs).strip()
            if not text:
                continue

            # Match by content heuristics after replacement
            if breadcrumb_val and text == breadcrumb_val.strip():
                for run in para.runs:
                    _apply_run_font(run, FONT_BREADCRUMB, bold=True, color=HPE_GREEN)

            elif module_name_val and text == module_name_val.strip():
                for run in para.runs:
                    _apply_run_font(run, FONT_MODULE_NAME, bold=True, color=TEXT_DARK)

            elif module_body_val and text == module_body_val.strip():
                for run in para.runs:
                    _apply_run_font(run, FONT_MODULE_BODY, bold=False, color=TEXT_DARK)


# ─── HTML → line blocks ───────────────────────────────────────────────────────

def _cell_text(td) -> str:
    return " ".join(td.get_text(separator=" ").split())


def html_to_blocks(raw_html: str):
    """
    Convert HTML to a list of blocks:
      {"type": "text",  "lines": [{"text": str, "depth": int, "kind": str}]}
      {"type": "table", "rows": [[str, ...], ...]}

    kind: "normal" | "bullet" | "numbered" | "heading" | "blank"
    depth: nesting level (0=top, 1=nested, ...)
    numbered lines also carry "counter": int
    """
    if not raw_html:
        return []

    if not HAS_BS4:
        text = _strip_html_plain(raw_html)
        if not text:
            return []
        lines = [{"text": ln, "depth": 0, "kind": "normal"}
                 for ln in text.split("\n") if ln.strip()]
        return [{"type": "text", "lines": lines}]

    soup   = BeautifulSoup(raw_html, "html.parser")
    blocks = []
    lines  = []

    def flush():
        content = [l for l in lines if l.get("text", "").strip() or l.get("kind") == "blank"]
        if any(l.get("text", "").strip() for l in content):
            blocks.append({"type": "text", "lines": list(content)})
        lines.clear()

    def norm(s):
        return html_mod.unescape(" ".join(str(s).split()))

    def parse_node(el, depth=0):
        tag = getattr(el, "name", None)

        if isinstance(el, NavigableString):
            t = norm(el)
            if t:
                lines.append({"text": t, "depth": depth, "kind": "normal"})
            return

        if tag in ("script", "style"):
            return

        if tag == "br":
            lines.append({"text": "", "depth": 0, "kind": "blank"})
            return

        if tag == "table":
            flush()
            rows = []
            for tr in el.find_all("tr"):
                cells = [_cell_text(td) for td in tr.find_all(["td", "th"])]
                if any(cells):
                    rows.append(cells)
            if rows:
                blocks.append({"type": "table", "rows": rows})
            return

        if tag in ("ul", "ol"):
            lines.append({"text": "", "depth": 0, "kind": "blank"})
            for idx, li in enumerate(el.find_all("li", recursive=False), start=1):
                direct_parts = []
                sub_lists    = []
                for child in li.children:
                    child_tag = getattr(child, "name", None)
                    if child_tag in ("ul", "ol"):
                        sub_lists.append(child)
                    elif isinstance(child, NavigableString):
                        t = norm(child)
                        if t:
                            direct_parts.append(t)
                    else:
                        t = norm(child.get_text(separator=" ", strip=True))
                        if t:
                            direct_parts.append(t)

                item_text = " ".join(direct_parts).strip()
                if item_text:
                    if tag == "ol":
                        lines.append({
                            "text": item_text, "depth": depth,
                            "kind": "numbered", "counter": idx,
                        })
                    else:
                        lines.append({
                            "text": item_text, "depth": depth, "kind": "bullet",
                        })

                for sub in sub_lists:
                    parse_node(sub, depth=depth + 1)

            lines.append({"text": "", "depth": 0, "kind": "blank"})
            return

        if tag in ("p", "div", "section", "article", "blockquote"):
            lines.append({"text": "", "depth": 0, "kind": "blank"})
            for child in el.children:
                parse_node(child, depth)
            lines.append({"text": "", "depth": 0, "kind": "blank"})
            return

        if tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
            lines.append({"text": "", "depth": 0, "kind": "blank"})
            t = norm(el.get_text(separator=" ", strip=True))
            if t:
                lines.append({"text": t, "depth": depth, "kind": "heading"})
            lines.append({"text": "", "depth": 0, "kind": "blank"})
            return

        if tag in ("strong", "b", "em", "i", "span", "a", "label", "u", "code"):
            t = norm(el.get_text(separator=" ", strip=True))
            if t:
                lines.append({"text": t, "depth": depth, "kind": "normal"})
            return

        for child in el.children:
            parse_node(child, depth)

    for top_child in soup.children:
        parse_node(top_child, depth=0)

    flush()

    # Collapse consecutive blank lines; strip leading/trailing blanks
    cleaned = []
    for blk in blocks:
        if blk["type"] != "text":
            cleaned.append(blk)
            continue
        result     = []
        prev_blank = False
        for ln in blk["lines"]:
            is_blank = (ln.get("kind") == "blank" or not ln.get("text", "").strip())
            if is_blank:
                if not prev_blank:
                    result.append({"text": "", "depth": 0, "kind": "blank"})
                prev_blank = True
            else:
                result.append(ln)
                prev_blank = False
        while result and not result[0].get("text", "").strip():
            result.pop(0)
        while result and not result[-1].get("text", "").strip():
            result.pop()
        if result:
            cleaned.append({"type": "text", "lines": result})

    return cleaned


def _strip_html_plain(raw: str) -> str:
    if not raw:
        return ""
    text = re.sub(r"<br\s*/?>",        "\n", raw,  flags=re.I)
    text = re.sub(r"</p\s*>",          "\n", text, flags=re.I)
    text = re.sub(r"</li\s*>",         "\n", text, flags=re.I)
    text = re.sub(r"<li[^>]*>",        "\u2022 ", text, flags=re.I)
    text = re.sub(r"</?(ul|ol)[^>]*>", "\n", text, flags=re.I)
    text = re.sub(r"<[^>]+>",          "",   text)
    text = html_mod.unescape(text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


# ─── Line height estimator ────────────────────────────────────────────────────

def _estimate_lines(line_objects: list) -> int:
    total = 0
    for ln in line_objects:
        if ln.get("kind") == "blank" or not ln.get("text", "").strip():
            total += 1
            continue
        depth_penalty = ln.get("depth", 0) * 8
        effective     = max(80, CHARS_PER_LINE - depth_penalty)
        total += max(1, (len(ln["text"]) + effective - 1) // effective)
    return total


def _lines_height_emu(line_objects: list) -> int:
    return int(_estimate_lines(line_objects) * int(LINE_H)) + int(TEXT_GAP)


def _fits(current_y: int, needed_emu: int) -> bool:
    return (current_y + needed_emu) <= int(CONTENT_BOTTOM)


# ─── Slide cloning ────────────────────────────────────────────────────────────

def clone_slide(prs: Presentation, src_slide):
    new_slide = prs.slides.add_slide(src_slide.slide_layout)
    src_sp = src_slide.shapes._spTree
    dst_sp = new_slide.shapes._spTree
    for child in list(dst_sp):
        dst_sp.remove(child)
    for child in src_sp:
        dst_sp.append(copy.deepcopy(child))
    src_bg = src_slide._element.find(f"{{{PML}}}bg")
    if src_bg is not None:
        dst_el = new_slide._element
        old_bg = dst_el.find(f"{{{PML}}}bg")
        if old_bg is not None:
            dst_el.remove(old_bg)
        dst_el.insert(0, copy.deepcopy(src_bg))
    new_slide.name   = src_slide.name
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return new_slide


# ─── Placeholder utilities ────────────────────────────────────────────────────

def _get_placeholder(slide, ph_idx: int):
    for shape in slide.shapes:
        try:
            if (shape.has_text_frame and
                    shape.placeholder_format is not None and
                    shape.placeholder_format.idx == ph_idx):
                return shape
        except Exception:
            pass
    return None


def _remove_placeholder(slide, ph_idx: int):
    spTree = slide.shapes._spTree
    to_rm  = []
    for shape in slide.shapes:
        try:
            if (shape.placeholder_format is not None and
                    shape.placeholder_format.idx == ph_idx):
                to_rm.append(shape._element)
        except Exception:
            pass
    for el in to_rm:
        spTree.remove(el)


def _fit_text_in_body_placeholder(slide, ph_idx: int = 1,
                                    min_font: Pt = Pt(7),
                                    target_font: Pt = FONT_BODY):
    shape = _get_placeholder(slide, ph_idx)
    if shape is None:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    try:
        txBody = tf._txBody
        bodyPr = txBody.find(qn("a:bodyPr"))
        if bodyPr is not None:
            for tag in ("a:spAutoFit", "a:normAutoFit", "a:noAutofit"):
                for el in bodyPr.findall(qn(tag)):
                    bodyPr.remove(el)
            etree.SubElement(bodyPr, qn("a:normAutoFit"))
    except Exception:
        pass


# ─── Text replacement (handles split runs) ───────────────────────────────────

def replace_in_paragraph(para, fields: dict):
    if not para.runs:
        return
    full = "".join(r.text for r in para.runs)
    replaced = full
    for k, v in fields.items():
        if not k.startswith("__"):
            replaced = replaced.replace(k, str(v) if v is not None else "")
    if replaced == full:
        return
    para.runs[0].text = replaced
    for run in para.runs[1:]:
        run.text = ""


def fill_shape(shape, fields: dict):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            replace_in_paragraph(para, fields)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    replace_in_paragraph(para, fields)
    if shape.shape_type == 6:
        try:
            for sub in shape.shapes:
                fill_shape(sub, fields)
        except Exception:
            pass


def fill_slide(slide, fields: dict):
    for shape in slide.shapes:
        fill_shape(shape, fields)


def fill_slide_with_typography(slide, fields: dict):
    """
    Fill token values AND enforce typography for header tokens.
    Call this instead of fill_slide() for content slides so that
    breadcrumb/moduleName/moduleBody get correct font sizes.
    """
    fill_slide(slide, fields)
    _style_header_placeholder(slide, fields)


# ─── Core textbox writer ──────────────────────────────────────────────────────

def _write_lines_to_textbox(slide, line_objects: list, top_emu: int) -> int:
    """
    Write line_objects into a new textbox.  Returns new Y (bottom + gap).
    Each line dict: {"text": str, "depth": int, "kind": str, "counter"?: int}
    """
    if not line_objects:
        return top_emu

    n_lines = _estimate_lines(line_objects)
    height  = int(n_lines * int(LINE_H)) + int(Inches(0.10))

    txBox = slide.shapes.add_textbox(
        int(CONTENT_LEFT), top_emu,
        int(CONTENT_WIDTH), height,
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, ln in enumerate(line_objects):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        kind     = ln.get("kind", "normal")
        depth    = ln.get("depth", 0)
        text     = ln.get("text", "")
        counter  = ln.get("counter", 1)
        is_blank = (kind == "blank" or not text.strip())

        run = para.add_run()

        if is_blank:
            run.text      = " "
            run.font.size = Pt(3)
        else:
            if kind == "numbered":
                run.text = f"{counter}. {text}"
            elif kind == "bullet":
                run.text = f"\u2022 {text}"
            else:
                run.text = text

            run.font.size      = FONT_BODY
            run.font.color.rgb = TEXT_DARK
            run.font.bold      = (kind == "heading")
            run.font.italic    = False

        try:
            pPr = para._p.get_or_add_pPr()

            # Line spacing — 115%
            lnSpc  = etree.SubElement(pPr, qn("a:lnSpc"))
            lnPct  = etree.SubElement(lnSpc, qn("a:spcPct"))
            lnPct.set("val", "115000")

            # Space after
            spcAft    = etree.SubElement(pPr, qn("a:spcAft"))
            spcAftPts = etree.SubElement(spcAft, qn("a:spcPts"))
            if is_blank:
                spcAftPts.set("val", "0")
            elif kind == "heading":
                spcAftPts.set("val", "180")
            elif kind in ("bullet", "numbered"):
                spcAftPts.set("val", "120")
            else:
                spcAftPts.set("val", "60")

            # Space before heading
            spcBef    = etree.SubElement(pPr, qn("a:spcBef"))
            spcBefPts = etree.SubElement(spcBef, qn("a:spcPts"))
            if kind == "heading" and i > 0:
                spcBefPts.set("val", "160")
            else:
                spcBefPts.set("val", "0")

            # Indent for list items
            if kind in ("bullet", "numbered"):
                margin_l = LIST_BASE_L + depth * LIST_DEPTH + LIST_HANG
                pPr.set("marL",   str(margin_l))
                pPr.set("indent", str(-LIST_HANG))
            elif depth > 0:
                pPr.set("marL", str(depth * LIST_DEPTH))

        except Exception:
            pass

    return top_emu + height + int(TEXT_GAP)


# ─── PPTX table builder ───────────────────────────────────────────────────────

def _set_cell_bg(cell, rgb: RGBColor):
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for tag in ("a:solidFill", "a:gradFill", "a:noFill"):
        for old in tcPr.findall(qn(tag)):
            tcPr.remove(old)
    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgbClr   = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", str(rgb).upper())


def _set_cell_border(cell, rgb: RGBColor = None, width_pt: float = 0.5):
    """Add subtle cell borders for industry-standard table appearance."""
    if rgb is None:
        rgb = RGBColor(0xCC, 0xCC, 0xCC)
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    w_emu = int(width_pt * 12700)   # 1pt = 12700 EMU
    color_val = str(rgb).upper()

    for border_tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        existing = tcPr.find(qn(border_tag))
        if existing is not None:
            tcPr.remove(existing)
        ln_el   = etree.SubElement(tcPr, qn(border_tag))
        ln_el.set("w", str(w_emu))
        ln_el.set("cap", "flat")
        ln_el.set("cmpd", "sng")
        solidFill = etree.SubElement(ln_el, qn("a:solidFill"))
        srgbClr   = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", color_val)
        etree.SubElement(ln_el, qn("a:prstDash")).set("val", "solid")


def _add_pptx_table(slide, rows: list, left=None, top=None,
                    width=None, first_row_is_header=True) -> int:
    if not rows:
        return int(top) if top is not None else int(CONTENT_TOP)

    n_cols = max(len(r) for r in rows)
    n_rows = len(rows)
    rows   = [list(r) + [""] * (n_cols - len(r)) for r in rows]

    left_emu  = int(left)  if left  is not None else int(TABLE_LEFT)
    top_emu   = int(top)   if top   is not None else int(CONTENT_TOP)
    width_emu = int(width) if width is not None else int(TABLE_W)

    row_h = int(TABLE_ROW_H)
    hdr_h = int(TABLE_HDR_H)

    row_heights  = [hdr_h if (ri == 0 and first_row_is_header) else row_h
                    for ri in range(n_rows)]
    exact_height = sum(row_heights)

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols, left_emu, top_emu, width_emu, exact_height
    )
    tbl = tbl_shape.table

    base_w    = width_emu // n_cols
    remainder = width_emu - base_w * n_cols
    for ci in range(n_cols):
        tbl.columns[ci].width = base_w + (remainder if ci == n_cols - 1 else 0)

    for ri, row_data in enumerate(rows):
        is_hdr = ri == 0 and first_row_is_header
        bg     = HPE_GREEN if is_hdr else (STRIPE_ODD if ri % 2 == 1 else STRIPE_EVN)
        for ci, val in enumerate(row_data[:n_cols]):
            cell = tbl.cell(ri, ci)
            cell.text = str(val or "")
            cell.text_frame.word_wrap = True
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = FONT_HDR if is_hdr else FONT_TABLE
                    if is_hdr:
                        run.font.bold      = True
                        run.font.color.rgb = WHITE
                    else:
                        run.font.color.rgb = TEXT_DARK
                        run.font.bold      = False
            _set_cell_bg(cell, bg)
            # Subtle borders on data rows for professional appearance
            if not is_hdr:
                _set_cell_border(cell, RGBColor(0xDD, 0xDD, 0xDD), 0.5)

    for ri, rh in enumerate(row_heights):
        tbl.rows[ri].height = rh

    # Patch graphic frame bounding box to exact height
    try:
        sp_el = tbl_shape._element
        pxfrm = sp_el.find(qn('p:xfrm'))
        if pxfrm is not None:
            aext = pxfrm.find(qn('a:ext'))
            if aext is not None:
                aext.set('cy', str(exact_height))
        for aext in sp_el.iter(qn('a:ext')):
            aext.set('cy', str(exact_height))
            break
    except Exception:
        pass

    return top_emu + exact_height


# ─── Slide removal ────────────────────────────────────────────────────────────

def remove_slide(prs, idx: int):
    sldIdLst = prs.slides._sldIdLst
    sldId    = sldIdLst[idx]
    rId      = sldId.get(f"{{{REL}}}id")
    sldIdLst.remove(sldId)
    if rId:
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass


# ─── Content slide factory ────────────────────────────────────────────────────

def _make_content_slide(prs, content_src,
                         breadcrumb: str, moduleNum: str, moduleName: str):
    """
    Clone template, fill header tokens with correct typography, REMOVE body placeholder.
    All body content added as textboxes after this call.
    """
    slide = clone_slide(prs, content_src)
    fields = {
        "{{breadcrumb}}":  breadcrumb,
        "{{moduleNum}}":   moduleNum,
        "{{moduleName}}":  moduleName,
        "{{moduleBody}}":  "",
    }
    # Fill tokens
    fill_slide(slide, fields)
    # Apply typography to header elements after replacement
    _style_header_placeholder(slide, fields)
    _remove_placeholder(slide, ph_idx=1)
    return slide


# ─── Pagination helpers ───────────────────────────────────────────────────────

def _paginate_lines(line_objects: list, max_lines: int) -> list:
    """
    Split line_objects into chunks that each fit within max_lines.
    Returns list of chunks (each chunk = list of line dicts).
    """
    chunks  = []
    current = []
    count   = 0

    for ln in line_objects:
        kind     = ln.get("kind", "normal")
        is_blank = (kind == "blank" or not ln.get("text", "").strip())

        if is_blank:
            cost = 1
        else:
            depth_penalty = ln.get("depth", 0) * 8
            effective     = max(80, CHARS_PER_LINE - depth_penalty)
            cost = max(1, (len(ln.get("text", "")) + effective - 1) // effective)

        if count + cost > max_lines and current:
            while current and not current[-1].get("text", "").strip():
                current.pop()
            if current:
                chunks.append(current)
            current = []
            count   = 0
            if is_blank:
                continue

        current.append(ln)
        count += cost

    while current and not current[-1].get("text", "").strip():
        current.pop()
    if current:
        chunks.append(current)

    return chunks if chunks else [[]]


def _paginate_table_rows(data_rows: list, max_rows: int):
    if not data_rows:
        return
    for i in range(0, len(data_rows), max_rows):
        chunk = data_rows[i:i + max_rows]
        if chunk:
            yield chunk


# ─── HTML-aware slide builder ─────────────────────────────────────────────────

def build_html_slide(prs, content_src,
                     breadcrumb: str, moduleNum: str, moduleName: str,
                     body_html: str) -> int:
    """
    Build one or more slides.  Body content ONLY via textboxes — consistent
    spacing/indent on every slide regardless of template placeholder styles.
    Content overflow is strictly guarded: a new slide is created whenever
    remaining space is insufficient.
    """
    blocks = html_to_blocks(body_html)
    log    = lambda *a: print(*a, file=sys.stderr, flush=True)
    log(f"[build_html_slide] '{moduleName}' -> {len(blocks)} blocks")

    if not blocks:
        log(f"[build_html_slide] '{moduleName}' -> no content blocks, skipping slide")
        return 0

    slides_created  = 0
    cur_slide       = None
    cur_y           = int(CONTENT_TOP)
    slide_is_fresh  = False

    def ensure_slide():
        nonlocal cur_slide, cur_y, slides_created, slide_is_fresh
        if cur_slide is None:
            cur_slide      = _make_content_slide(
                prs, content_src, breadcrumb, moduleNum, moduleName
            )
            cur_y          = int(CONTENT_TOP)
            slides_created += 1
            slide_is_fresh  = True

    def new_slide():
        nonlocal cur_slide, cur_y, slides_created, slide_is_fresh
        cur_slide      = _make_content_slide(
            prs, content_src, breadcrumb, moduleNum, moduleName
        )
        cur_y          = int(CONTENT_TOP)
        slides_created += 1
        slide_is_fresh  = True

    def need_new_slide(needed_emu: int) -> bool:
        """
        Return True only when content genuinely won't fit AND the current
        slide already has some content.  Never break to a new slide when
        the slide was just created — that would leave it blank.
        """
        if slide_is_fresh:
            return False
        return not _fits(cur_y, needed_emu)

    for blk in blocks:

        if blk["type"] == "text":
            line_chunks = _paginate_lines(blk["lines"], MAX_LINES_PER_SLIDE)
            for chunk in line_chunks:
                if not chunk:
                    continue
                needed = _lines_height_emu(chunk)
                log(f"[build_html_slide] text chunk {len(chunk)} lines, "
                    f"needed={needed/914400:.3f}in cur_y={cur_y/914400:.3f}in "
                    f"bottom={int(CONTENT_BOTTOM)/914400:.3f}in fresh={slide_is_fresh}")
                ensure_slide()
                if need_new_slide(needed):
                    new_slide()
                    log(f"[build_html_slide] text overflow -> slide {slides_created}")
                cur_y          = _write_lines_to_textbox(cur_slide, chunk, cur_y)
                slide_is_fresh = False

        elif blk["type"] == "table":
            rows = blk["rows"]
            if not rows:
                continue
            header_row = rows[0]
            data_rows  = rows[1:]
            if not data_rows:
                data_rows = [[]]
            for chunk_data in _paginate_table_rows(data_rows, MAX_TABLE_ROWS_PER_SLIDE):
                chunk_rows = [header_row] + chunk_data
                n_chunk    = len(chunk_rows)
                tbl_h      = int(TABLE_HDR_H) + (n_chunk - 1) * int(TABLE_ROW_H)
                log(f"[build_html_slide] table {n_chunk} rows, "
                    f"tbl_h={tbl_h/914400:.3f}in cur_y={cur_y/914400:.3f}in "
                    f"fresh={slide_is_fresh}")
                ensure_slide()
                if need_new_slide(tbl_h):
                    new_slide()
                    log(f"[build_html_slide] table overflow -> slide {slides_created}")
                log(f"[build_html_slide] table -> slide {slides_created}")
                bottom = _add_pptx_table(
                    cur_slide, chunk_rows,
                    left=int(TABLE_LEFT), top=cur_y,
                    width=int(TABLE_W), first_row_is_header=True,
                )
                cur_y          = bottom + int(TEXT_GAP)
                slide_is_fresh = False

    log(f"[build_html_slide] -> {slides_created} slides total")
    return max(slides_created, 1)


# ─── Main ─────────────────────────────────────────────────────────────────────

def main():
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    if hasattr(sys.stderr, "reconfigure"):
        sys.stderr.reconfigure(encoding="utf-8", errors="replace")

    if len(sys.argv) < 2:
        print("Usage: fill_proposal.py <data.json>", file=sys.stderr)
        sys.exit(1)

    with open(sys.argv[1], encoding="utf-8") as f:
        data = json.load(f)

    prs        = Presentation(data["template"])
    n_template = len(prs.slides)

    cover_src = content_src = closing_src = None
    for slide in prs.slides:
        lname = slide.slide_layout.name
        if "Cover"  in lname and cover_src   is None:
            cover_src   = slide
        elif "Thank" in lname and closing_src is None:
            closing_src = slide
        elif content_src is None:
            content_src = slide

    if cover_src   is None: cover_src   = prs.slides[0]
    if content_src is None: content_src = prs.slides[1]
    if closing_src is None: closing_src = prs.slides[min(2, n_template - 1)]

    generated = 0

    for spec in data["slides"]:
        layout = spec["layout"]
        fields = spec.get("fields", {})

        if layout == "LAYOUT_COVER":
            slide = clone_slide(prs, cover_src)
            fill_slide(slide, fields)
            generated += 1

        elif layout == "LAYOUT_CLOSING":
            slide = clone_slide(prs, closing_src)
            fill_slide(slide, fields)
            generated += 1

        elif layout == "LAYOUT_CONTENT":
            slide = clone_slide(prs, content_src)
            ph = _get_placeholder(slide, ph_idx=1)
            if ph:
                ph.left   = int(CONTENT_LEFT)
                ph.top    = int(CONTENT_TOP)
                ph.width  = int(CONTENT_WIDTH)
                ph.height = int(CONTENT_BOTTOM) - int(CONTENT_TOP)
            # Fill tokens then apply typography
            fill_slide(slide, fields)
            _style_header_placeholder(slide, fields)
            _fit_text_in_body_placeholder(slide, ph_idx=1,
                                          min_font=Pt(7), target_font=FONT_BODY)
            generated += 1

        elif layout == "LAYOUT_HTML":
            breadcrumb = fields.get("{{breadcrumb}}", "")
            moduleNum  = fields.get("{{moduleNum}}",  "")
            moduleName = fields.get("{{moduleName}}", "")
            body_html  = fields.get("__html__",       "")
            n = build_html_slide(
                prs, content_src,
                breadcrumb, moduleNum, moduleName, body_html,
            )
            generated += n

        elif layout == "LAYOUT_TABLE":
            text_fields = {k: v for k, v in fields.items()
                           if not k.startswith("__")}
            text_fields.setdefault("{{moduleNum}}",  "")
            text_fields.setdefault("{{moduleName}}", "")
            text_fields.setdefault("{{moduleBody}}", "")

            headers  = fields.get("__tableHeader__") or []
            all_rows = fields.get("__tableRows__")    or []

            if not all_rows:
                slide = clone_slide(prs, content_src)
                fill_slide(slide, text_fields)
                _style_header_placeholder(slide, text_fields)
                _remove_placeholder(slide, ph_idx=1)
                generated += 1
            else:
                for chunk_idx, chunk in enumerate(
                    _paginate_table_rows(all_rows, MAX_TABLE_ROWS_PER_SLIDE)
                ):
                    slide = clone_slide(prs, content_src)
                    chunk_fields = dict(text_fields)
                    if chunk_idx > 0:
                        chunk_fields["{{breadcrumb}}"] = (
                            text_fields.get("{{breadcrumb}}", "") + " (cont.)"
                        )
                    fill_slide(slide, chunk_fields)
                    _style_header_placeholder(slide, chunk_fields)
                    _remove_placeholder(slide, ph_idx=1)
                    table_rows = ([headers] + chunk) if headers else chunk
                    _add_pptx_table(
                        slide, table_rows,
                        left=int(TABLE_LEFT),
                        top=int(CONTENT_TOP),
                        width=int(TABLE_W),
                        first_row_is_header=bool(headers),
                    )
                    generated += 1

    for i in range(n_template - 1, -1, -1):
        remove_slide(prs, i)

    prs.save(data["output"])
    print(f"[SUCCESS] {generated} slides -> {data['output']}")


if __name__ == "__main__":
    main()